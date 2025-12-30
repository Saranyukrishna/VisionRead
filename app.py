import streamlit as st
import fitz
from docx import Document
from pptx import Presentation
from PIL import Image, ImageEnhance
import io
import os
from pathlib import Path
import base64
import numpy as np
import cohere
import google.generativeai as genai
from dotenv import load_dotenv
import tempfile
import shutil
from langchain.messages import HumanMessage, AIMessage
from tavily import TavilyClient
from groq import Groq
import time

# Load environment variables first
load_dotenv()

# Initialize API clients with error handling
@st.cache_resource(show_spinner=False)
def initialize_clients():
    try:
        cohere_api_key = os.getenv("COHERE_API_KEY")
        gemini_api_key = os.getenv("GOOGLE_API_KEY")
        tavily_api_key = os.getenv("TAVILY_API_KEY")
        groq_api_key = os.getenv("GROQ_API_KEY")

        if not all([cohere_api_key, gemini_api_key, tavily_api_key, groq_api_key]):
            st.error("API keys not found. Please check your .env file")
            st.stop()

        co = cohere.Client(api_key=cohere_api_key, timeout=30)
        genai.configure(api_key=gemini_api_key)
        tavily = TavilyClient(api_key=tavily_api_key)
        groq = Groq(api_key=groq_api_key, timeout=30)
        
        return co, genai, tavily, groq
    except Exception as e:
        st.error(f"Failed to initialize API clients: {str(e)}")
        st.stop()

co, genai, tavily, groq = initialize_clients()

# Constants
MAX_PIXELS = 1568 * 1568
SUPPORTED_TYPES = ["pdf", "docx", "pptx"]
GEMINI_MODEL = "gemini-1.5-flash-latest"
IMAGE_QUALITY = 95
BLANK_IMAGE_THRESHOLD = 0.95

# Setup directories
@st.cache_resource(show_spinner=False)
def setup_directories():
    OUTPUT_DIR = Path(tempfile.mkdtemp())
    IMAGES_DIR = OUTPUT_DIR / "images"
    IMAGES_DIR.mkdir(parents=True, exist_ok=True)  # Ensure images directory exists
    TEXT_FILE = OUTPUT_DIR / "extracted_text.txt"
    return OUTPUT_DIR, IMAGES_DIR, TEXT_FILE

OUTPUT_DIR, IMAGES_DIR, TEXT_FILE = setup_directories()

# Cleanup function
def cleanup():
    if OUTPUT_DIR.exists():
        try:
            shutil.rmtree(OUTPUT_DIR)
            # Recreate the directory structure
            OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
            IMAGES_DIR.mkdir(parents=True, exist_ok=True)
        except Exception as e:
            st.warning(f"Cleanup warning: {str(e)}")

# Image processing functions
def is_blank_image(pil_image, threshold=BLANK_IMAGE_THRESHOLD):
    try:
        if pil_image.mode != 'RGB':
            pil_image = pil_image.convert('RGB')
        img_array = np.array(pil_image)
        dark_pixels = np.sum((img_array[:,:,0] < 50) &
                      (img_array[:,:,1] < 50) &
                      (img_array[:,:,2] < 50))
        total_pixels = img_array.shape[0] * img_array.shape[1]
        dark_ratio = dark_pixels / total_pixels

        white_pixels = np.sum((img_array[:,:,0] > 200) & 
                         (img_array[:,:,1] > 200) & 
                         (img_array[:,:,2] > 200))
        white_ratio = white_pixels / total_pixels

        return dark_ratio > threshold or white_ratio > threshold
    except Exception as e:
        st.warning(f"Blank image check failed: {str(e)}")
        return False

def save_image(_image_pil, image_count):
    try:
        if not IMAGES_DIR.exists():
            IMAGES_DIR.mkdir(parents=True, exist_ok=True)
            
        img_path = IMAGES_DIR / f"image_{image_count:04d}.png"
        
        if _image_pil.mode in ('RGBA', 'LA', 'P'):
            _image_pil = _image_pil.convert('RGB')
        
        _image_pil.save(
            img_path,
            format="PNG",
            quality=IMAGE_QUALITY,
            optimize=True,
            compress_level=6
        )
        
        return str(img_path.resolve())
    
    except Exception as e:
        st.error(f"Error saving image {image_count}: {str(e)}")
        return None

def resize_image(pil_image, max_pixels=MAX_PIXELS):
    try:
        org_width, org_height = pil_image.size
        if org_width * org_height > max_pixels:
            scale_factor = (max_pixels / (org_width * org_height)) ** 0.5
            new_width = int(org_width * scale_factor)
            new_height = int(org_height * scale_factor)
            pil_image = pil_image.resize((new_width, new_height), Image.Resampling.LANCZOS)
        return pil_image
    except Exception as e:
        st.warning(f"Image resize failed: {str(e)}")
        return pil_image

def base64_from_image(img_path):
    try:
        if not Path(img_path).exists():
            return None
            
        pil_image = Image.open(img_path)
        img_format = pil_image.format or "PNG"
        if pil_image.mode in ('RGBA', 'LA'):
            pil_image = pil_image.convert('RGB')
        pil_image = resize_image(pil_image)
        with io.BytesIO() as buffer:
            pil_image.save(buffer, format=img_format, quality=IMAGE_QUALITY, optimize=True)
            encoded = base64.b64encode(buffer.getvalue()).decode("utf-8")
            return f"data:image/{img_format.lower()};base64,{encoded}"
    except Exception as e:
        st.error(f"Error processing image: {str(e)}")
        return None

# Document processing functions
def process_document(uploaded_file):
    text = ""
    image_paths = []
    image_count = 1
    
    try:
        file_ext = uploaded_file.name.split(".")[-1].lower()
        
        if file_ext == "pdf":
            with fitz.open(stream=uploaded_file.read(), filetype="pdf") as pdf:
                for page in pdf:
                    text += page.get_text()
                    for img in page.get_images(full=True):
                        xref = img[0]
                        base_image = pdf.extract_image(xref)
                        img_bytes = base_image["image"]
                        img_pil = Image.open(io.BytesIO(img_bytes))
                        if not is_blank_image(img_pil):
                            img_path = save_image(img_pil, image_count)
                            if img_path:
                                image_paths.append(img_path)
                                image_count += 1
        
        elif file_ext == "docx":
            doc = Document(uploaded_file)
            for para in doc.paragraphs:
                text += para.text + "\n"
            for rel in doc.part._rels:
                rel_obj = doc.part._rels[rel]
                if "image" in rel_obj.target_ref:
                    img_data = rel_obj.target_part.blob
                    img_pil = Image.open(io.BytesIO(img_data))
                    if not is_blank_image(img_pil):
                        img_path = save_image(img_pil, image_count)
                        if img_path:
                            image_paths.append(img_path)
                            image_count += 1
        
        elif file_ext == "pptx":
            prs = Presentation(uploaded_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                    if shape.shape_type == 13:
                        img_stream = shape.image.blob
                        img_pil = Image.open(io.BytesIO(img_stream))
                        if not is_blank_image(img_pil):
                            img_path = save_image(img_pil, image_count)
                            if img_path:
                                image_paths.append(img_path)
                                image_count += 1
        
        # Save text to file
        with open(TEXT_FILE, "w", encoding="utf-8") as f:
            f.write(text)
            
        return text, image_paths
    
    except Exception as e:
        st.error(f"Failed to process document: {str(e)}")
        cleanup()
        return None, None

# AI Functions with error handling and timeouts
def ask_gemini(question, context=None, img_path=None):
    try:
        model = genai.GenerativeModel(GEMINI_MODEL)
        start_time = time.time()
        
        if img_path and context:
            prompt = f"""Use the context and image to answer the question if relevant.
Context: {context}
Question: {question}"""
            img = Image.open(img_path)
            response = model.generate_content([prompt, img], request_options={"timeout": 30})
        elif img_path:
            prompt = f"""Analyze the image to answer the question if relevant.
Question: {question}"""
            img = Image.open(img_path)
            response = model.generate_content([prompt, img], request_options={"timeout": 30})
        elif context:
            prompt = f"""Use the context to answer the question if relevant.
Context: {context}
Question: {question}"""
            response = model.generate_content(prompt, request_options={"timeout": 30})
        else:
            response = model.generate_content(question, request_options={"timeout": 30})
            
        if time.time() - start_time > 25:
            st.warning("Gemini response took longer than expected")
            
        return response.text if response else "No response from Gemini"
    
    except Exception as e:
        st.error(f"Gemini error: {str(e)}")
        return f"Error querying Gemini: {str(e)}"

def search_tavily(query, search_depth='advanced', max_results=5):
    try:
        start_time = time.time()
        response = tavily.search(
            query=query,
            include_answer=True,
            include_raw_content=True,
            include_sources=True,
            max_results=max_results,
            search_depth=search_depth,
            timeout=30
        )
        if time.time() - start_time > 25:
            st.warning("Web search took longer than expected")
        return response
    except Exception as e:
        st.error(f"Tavily search error: {str(e)}")
        return None

def ask_groq(question, context=None):
    try:
        messages = []
        if context:
            messages.append({
                "role": "system",
                "content": f"Use this context if relevant: {context}"
            })
        messages.append({
            "role": "user",
            "content": question
        })

        start_time = time.time()
        response = groq.chat.completions.create(
            model="llama-3.1-8b-instant",
            messages=messages,
            temperature=0.8,
            timeout=30
        )
        
        if time.time() - start_time > 25:
            st.warning("Groq response took longer than expected")
            
        return response.choices[0].message.content if response else "No response from Groq"
    
    except Exception as e:
        st.error(f"Groq error: {str(e)}")
        return f"Error querying Groq: {str(e)}"

# UI Rendering Functions
def render_chat(container, chat_history):
    with container:
        for message in chat_history:
            if isinstance(message, HumanMessage):
                st.markdown(
                    f"<div style='text-align: right; color: white; background-color: #0a84ff; padding: 8px; border-radius: 10px; margin: 5px 0; max-width: 80%; float: right; clear: both;'>{message.content}</div>",
                    unsafe_allow_html=True
                )
            elif isinstance(message, AIMessage):
                st.markdown(
                    f"<div style='text-align: left; color: black; background-color: #d1d1d1; padding: 8px; border-radius: 10px; margin: 5px 0; max-width: 80%; float: left; clear: both;'>{message.content}</div>",
                    unsafe_allow_html=True
                )

# Initialize Streamlit app
st.set_page_config(page_title="Document Q&A", layout="wide")
st.title("📄 Document Q&A with Image Analysis")

# Initialize session state
if 'text_chat_history' not in st.session_state:
    st.session_state.text_chat_history = []
if 'image_chat_history' not in st.session_state:
    st.session_state.image_chat_history = []
if 'general_chat_history' not in st.session_state:
    st.session_state.general_chat_history = []
if 'processed' not in st.session_state:
    st.session_state.processed = False
    st.session_state.text = ""
    st.session_state.image_paths = []
    st.session_state.selected_img = None
if 'prev_uploaded_file' not in st.session_state:
    st.session_state.prev_uploaded_file = None

# File upload section
with st.expander("Upload Document", expanded=True):
    uploaded_file = st.file_uploader("Choose a file", type=SUPPORTED_TYPES, key="file_uploader")
    
    # Clear state when file is removed
    if not uploaded_file and st.session_state.prev_uploaded_file:
        st.session_state.processed = False
        st.session_state.text = ""
        st.session_state.image_paths = []
        st.session_state.selected_img = None
        st.session_state.prev_uploaded_file = None
        st.session_state.text_chat_history = []
        st.session_state.image_chat_history = []
        cleanup()
        st.rerun()

# Process document when uploaded
if uploaded_file and uploaded_file != st.session_state.prev_uploaded_file:
    with st.spinner("Processing document..."):
        # Clear previous state
        st.session_state.text = ""
        st.session_state.image_paths = []
        st.session_state.selected_img = None
        cleanup()
        
        # Process new document
        st.session_state.text, st.session_state.image_paths = process_document(uploaded_file)
        if st.session_state.text is not None:
            st.session_state.processed = True
            st.session_state.prev_uploaded_file = uploaded_file
            st.success("Document processed successfully!")
        else:
            st.session_state.processed = False

# Tabs interface
tab1, tab2, tab3 = st.tabs(["📝 Text Analysis", "🖼️ Image Analysis", "💬 General Chat"])

with tab1:
    st.subheader("Text Analysis")
    
    if not st.session_state.processed:
        st.info("Please upload a document to analyze")
    else:
        # Initialize session state variables
        if "text_chat_history" not in st.session_state:
            st.session_state.text_chat_history = []
        if "last_text_question" not in st.session_state:
            st.session_state.last_text_question = ""

        # Clear chat option
        if st.button("Clear Chat", key="clear_text_chat"):
            st.session_state.text_chat_history = []
            st.session_state.last_text_question = ""
            st.rerun()

        # Display extracted text
        with st.expander("View Extracted Text"):
            st.text_area("Extracted Text", st.session_state.text, height=200, label_visibility="collapsed")

        # Chat history container
        text_chat_container = st.container(height=400)
        render_chat(text_chat_container, st.session_state.text_chat_history)

        # User chat input
        user_text_input = st.chat_input("Ask about the text content...")

        # Handle new question
        if user_text_input and user_text_input != st.session_state.last_text_question:
            st.session_state.last_text_question = user_text_input
            st.session_state.text_chat_history.append(HumanMessage(content=user_text_input))

            with st.spinner("Analyzing text..."):
                try:
                    context = "\n".join(
                        f"{'User' if isinstance(m, HumanMessage) else 'Assistant'}: {m.content}"
                        for m in st.session_state.text_chat_history[-4:]
                    )

                    prompt = f"""Question: {user_text_input}

Extracted Text:
{st.session_state.text}

Chat Context:
{context if context else 'No previous chat context'}

Please answer the question using the text above. If the answer cannot be found, say so clearly."""

                    answer = ask_gemini(prompt)
                    st.session_state.text_chat_history.append(AIMessage(content=answer))
                    st.rerun()

                except Exception as e:
                    st.session_state.text_chat_history.append(
                        AIMessage(content=f"⚠️ Error: {str(e)}")
                    )
                    st.rerun()

with tab2:
    st.subheader("Image Analysis")
    
    if not st.session_state.processed:
        st.info("Please upload a document to analyze")
    else:
        img_col, _ = st.columns([1, 3])
        with img_col:
            if st.session_state.selected_img:
                try:
                    if Path(st.session_state.selected_img).exists():
                        selected_img = Image.open(st.session_state.selected_img)
                        
                        with st.expander("Image Enhancement Options"):
                            enhance = st.checkbox("Enhance Image Quality", value=True)
                            contrast = st.slider("Adjust Contrast", 0.5, 2.0, 1.0)
                            sharpness = st.slider("Adjust Sharpness", 0.0, 2.0, 1.0)
                            
                            if enhance:
                                enhancer = ImageEnhance.Contrast(selected_img)
                                selected_img = enhancer.enhance(contrast)
                                enhancer = ImageEnhance.Sharpness(selected_img)
                                selected_img = enhancer.enhance(sharpness)
                        st.image(selected_img, 
                                 caption="Selected Image", 
                                 use_container_width=True,
                                 output_format="PNG")
                        with io.BytesIO() as buffer:
                            selected_img.save(buffer, format="PNG", quality=IMAGE_QUALITY)
                            st.download_button(
                                label="Download Enhanced Image",
                                data=buffer.getvalue(),
                                file_name="enhanced_image.png",
                                mime="image/png"
                            )
                    else:
                        st.warning("Selected image no longer exists")
                        st.session_state.selected_img = None
                except Exception as e:
                    st.error(f"Error loading selected image: {str(e)}")
            else:
                st.info("Select an image from below to analyze")
        
        if st.session_state.selected_img and Path(st.session_state.selected_img).exists():
            image_chat_container = st.container()
            render_chat(image_chat_container, st.session_state.image_chat_history)
        
            # Input section at the bottom
            input_col = st.container()
            with input_col:
                st.write("**Ask about the image**")
                user_image_input = st.text_input(
                    "Ask about the image:", 
                    key="image_input", 
                    placeholder="Type your question here...",
                    label_visibility="collapsed",
                    disabled=not st.session_state.selected_img
                )
                image_send_button = st.button("Send", key="image_send", disabled=not st.session_state.selected_img)
            
            # Handle user input and generate responses
            if image_send_button and user_image_input:
                st.session_state.image_chat_history.append(HumanMessage(content=user_image_input))
                if user_image_input.lower() == 'close the chat':
                    st.stop()
                
                with st.spinner("Analyzing image..."):
                    answer = ask_gemini(
                        user_image_input, 
                        img_path=st.session_state.selected_img, 
                        context=st.session_state.text
                    )
                    st.session_state.image_chat_history.append(AIMessage(content=answer))
                    st.rerun()
        
        if st.session_state.image_paths:
            st.divider()
            st.write("Document Images:")
            num_cols = 4
            valid_image_paths = [p for p in st.session_state.image_paths if Path(p).exists()]
            rows = (len(valid_image_paths) + num_cols - 1) // num_cols
            
            for row in range(rows):
                cols = st.columns(num_cols)
                for col_idx in range(num_cols):
                    img_idx = row * num_cols + col_idx
                    if img_idx < len(valid_image_paths):
                        img_path = valid_image_paths[img_idx]
                        with cols[col_idx]:
                            try:
                                img = Image.open(img_path)
                                if not is_blank_image(img):
                                    st.image(img, use_container_width=True, output_format="PNG")
                                    if st.button(f"Select Image {img_idx+1}", key=f"btn_{img_idx}"):
                                        st.session_state.selected_img = img_path
                                        st.session_state.image_chat_history = []
                                        st.rerun()
                            except Exception as e:
                                st.error(f"Error loading image: {str(e)}")
            if not valid_image_paths:
                st.warning("No valid images found in the document")
        else:
            st.write("No images found in the document.")

with tab3:
    st.subheader("General Chat")

    # Initialize session state variables
    if "general_chat_history" not in st.session_state:
        st.session_state.general_chat_history = []
    if "last_question" not in st.session_state:
        st.session_state.last_question = ""

    # Configuration toggles
    col1, col2, col3 = st.columns(3)
    with col1:
        use_groq = st.toggle("Use Groq (ultra fast)", value=True, key="groq_toggle")
    with col2:
        enable_search = st.toggle("Enable web search", value=True, key="search_toggle")
    with col3:
        if st.button("Clear Chat", key="clear_chat"):
            st.session_state.general_chat_history = []
            st.session_state.last_question = ""
            st.rerun()

    # Display chat history
    chat_container = st.container(height=400)
    render_chat(chat_container, st.session_state.general_chat_history)

    # Chat input
    user_input = st.chat_input("Ask any question...")

    # Handle new question submission
    if user_input and user_input != st.session_state.last_question:
        st.session_state.last_question = user_input
        st.session_state.general_chat_history.append(HumanMessage(content=user_input))

        with st.spinner("Generating response..."):
            try:
                context = "\n".join(
                    f"{'User' if isinstance(m, HumanMessage) else 'Assistant'}: {m.content}"
                    for m in st.session_state.general_chat_history[-4:]
                )

                prompt = f"""New question: {user_input}

Previous conversation context:
{context if context else 'No previous context'}

Please provide a fresh, concise response to the new question above. 
If you need more information, say so explicitly."""

                initial_answer = ask_groq(prompt) if use_groq else ask_gemini(prompt)

                # Determine if web search is needed
                triggers = [
                    "i don't know", "not sure", "as of my knowledge",
                    "as of my last update", "i don't have information",
                    "my training data only goes up to", "i couldn't find"
                ]
                needs_search = enable_search and any(t in initial_answer.lower() for t in triggers)

                if needs_search:
                    st.session_state.general_chat_history.append(
                        AIMessage(content=initial_answer + "\n\n(Searching for updated information...)")
                    )
                    st.rerun()
                else:
                    st.session_state.general_chat_history.append(AIMessage(content=initial_answer))
                    st.rerun()

            except Exception as e:
                st.session_state.general_chat_history.append(
                    AIMessage(content=f"⚠️ Error: {str(e)}")
                )
                st.rerun()

    # Search flow continuation
    if (
        st.session_state.general_chat_history
        and isinstance(st.session_state.general_chat_history[-1], AIMessage)
        and "(Searching for updated information...)" in st.session_state.general_chat_history[-1].content
    ):
        user_input = st.session_state.last_question
        try:
            search_results = search_tavily(user_input)
            initial_answer = st.session_state.general_chat_history[-1].content.replace("\n\n(Searching for updated information...)", "")

            if search_results and search_results.get('results'):
                relevant_links = "\n".join(
                    f"- [{r['title']}]({r['url']})" for r in search_results['results'][:3]
                )

                enhancement_prompt = f"""Original question: {user_input}

Initial response: {initial_answer}

New information from web search:
{search_results.get('answer', 'No summary available')}

Available sources:
{relevant_links}

Please provide an improved answer incorporating this new information when relevant.
Always cite sources using markdown links when using specific information."""

                final_answer = ask_groq(enhancement_prompt) if use_groq else ask_gemini(enhancement_prompt)

                final_output = f"{initial_answer}\n\n---\n\n**Updated Information**:\n{final_answer}"
            else:
                final_output = f"{initial_answer}\n\n(Web search didn't find additional information)"

            st.session_state.general_chat_history[-1] = AIMessage(content=final_output)
            st.rerun()

        except Exception as e:
            st.session_state.general_chat_history[-1] = AIMessage(
                content=f"⚠️ Error during Tavily search: {str(e)}"
            )
            st.rerun()

# Cleanup on app exit
st.session_state.cleanup = cleanup
